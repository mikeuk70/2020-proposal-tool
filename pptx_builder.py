"""
20.20 Proposal PPTX Builder — python-pptx version 2
Clean builds matching real 20.20 proposal layouts.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re

W = Inches(13.33)
H = Inches(7.5)

WHITE  = RGBColor(0xFF,0xFF,0xFF)
BLACK  = RGBColor(0x1A,0x1A,0x1A)
DARK   = RGBColor(0x11,0x14,0x18)
MID    = RGBColor(0x2D,0x2D,0x2D)
GREY   = RGBColor(0x88,0x88,0x88)
LGREY  = RGBColor(0xF2,0xF1,0xEE)
RULE   = RGBColor(0xDD,0xDB,0xD5)
DEFAULT_ACCENT = RGBColor(0xE9,0x71,0x32)

F_HEAD = 'Filson Pro Heavy'
F_BODY = 'Filson Pro'

CLUB_COLOURS = {
    'aston villa':(0x5C,0x1A,0x2E),'villa':(0x5C,0x1A,0x2E),
    'newcastle':(0xC9,0xA8,0x4C),'nufc':(0xC9,0xA8,0x4C),
    'brighton':(0x00,0x57,0xB8),'bhafc':(0x00,0x57,0xB8),
    'arsenal':(0xEF,0x01,0x07),'liverpool':(0xC8,0x10,0x2E),
    'chelsea':(0x03,0x46,0x94),'crystal palace':(0x1B,0x45,0x8F),
    'cpfc':(0x1B,0x45,0x8F),'leeds':(0xFF,0xCD,0x00),'lufc':(0xFF,0xCD,0x00),
    'sunderland':(0xEB,0x17,0x2B),'safc':(0xEB,0x17,0x2B),
    'west ham':(0x7A,0x26,0x3A),'manchester city':(0x6C,0xAB,0xDD),
    'man city':(0x6C,0xAB,0xDD),'manchester united':(0xDA,0x29,0x1C),
    'man utd':(0xDA,0x29,0x1C),'tottenham':(0x13,0x22,0x57),
    'spurs':(0x13,0x22,0x57),'luton':(0xF7,0x83,0x1A),
    'luton town':(0xF7,0x83,0x1A),'sheffield':(0xEE,0x27,0x37),
    'nottingham':(0xE5,0x32,0x33),'wolves':(0xFD,0xB9,0x13),
    'wolverhampton':(0xFD,0xB9,0x13),'everton':(0x00,0x33,0x99),
    'celtic':(0x00,0x84,0x3D),'rangers':(0x00,0x33,0xA0),
}

def _accent(name):
    cl = (name or '').lower()
    for k,v in CLUB_COLOURS.items():
        if k in cl: return RGBColor(*v)
    return DEFAULT_ACCENT

def _clean(t):
    if not t: return ''
    # Strip markdown bold/italic
    t = re.sub(r'\*\*([^*]+)\*\*', r'\1', t)
    t = re.sub(r'\*([^*]+)\*', r'\1', t)
    # Strip markdown headers
    t = re.sub(r'^#{1,3}\s*', '', t, flags=re.MULTILINE)
    # Strip any remaining asterisks used as bold markers
    t = re.sub(r'\*\*|__', '', t)
    # Normalise excessive newlines
    t = re.sub(r'\n{3,}', '\n\n', t)
    return t.strip()

def _is_heading_line(s):
    """True if this line is a section label Claude added, not real content."""
    if len(s) > 60: return False
    if '.' in s or ',' in s: return False
    if re.search(r'stage \d|riba|objective|process|deliverable|meetings|presentations|your brief|cover letter|our approach|next steps|fees and timing', s, re.I):
        return True
    return False

def _parse(txt):
    """Return list of (type, text): 'prose', 'bullet', 'heading'."""
    out = []
    for line in _clean(txt).split('\n'):
        s = line.strip()
        if not s: continue
        if s.startswith(('-','\u2022','*')) or re.match(r'^\d+[.):]',s):
            item = re.sub(r'^[-\u2022*]\s*|^\d+[.):]\s*','',s)
            if len(item) > 5: out.append(('bullet', item))
        elif _is_heading_line(s):
            out.append(('heading', s))
        else:
            out.append(('prose', s))
    return out

def _prose(txt, n=4):
    parts = [c for t,c in _parse(txt) if t=='prose']
    # Drop any line that is clearly a stage/phase header Claude added
    parts = [p for p in parts if not re.search(
        r'(stage|phase)\s+\d.*?(riba|week|\|)', p, re.I)]
    text = ' '.join(parts)
    return ' '.join(re.split(r'(?<=[.!?])\s+', text)[:n]).strip()

def _bullets(txt, n=8):
    b = [c for t,c in _parse(txt) if t=='bullet']
    if not b:
        b = [c for t,c in _parse(txt) if t=='prose' and len(c)>20]
    return b[:n]

def _all_bullets(txt):
    """Return ALL bullet points — no cap. Used for slides so nothing is lost."""
    b = [c for t,c in _parse(txt) if t=='bullet']
    if not b:
        b = [c for t,c in _parse(txt) if t=='prose' and len(c)>20]
    return b

def _section_of(txt, label):
    """Extract a named sub-section from structured text."""
    pattern = re.compile(
        r'(?:^|\n)\s*' + re.escape(label) + r'[:\s]*\n(.*?)(?=\n\s*(?:Objective|Process|Deliverables|Meetings|$))',
        re.IGNORECASE | re.DOTALL
    )
    m = pattern.search(txt)
    if m:
        return m.group(1).strip()
    return ''

# ── DRAWING PRIMITIVES ────────────────────────────────────────────────────────

def _bg(slide, colour):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour

def _box(slide, x, y, w, h, colour):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = colour
    s.line.fill.background()
    return s

def _rule(slide, x, y, w, colour=None):
    s = slide.shapes.add_shape(1, x, y, w, Pt(0.75))
    s.fill.solid(); s.fill.fore_color.rgb = colour or RULE
    s.line.fill.background()

def _logo(slide, dark=False):
    c = WHITE if dark else BLACK
    tb = slide.shapes.add_textbox(W-Inches(0.88), Inches(0.12), Inches(0.72), Inches(0.62))
    tf = tb.text_frame
    for i, txt in enumerate(['20','20']):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = txt
        r.font.name = F_HEAD; r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = c
    dot = slide.shapes.add_shape(1, W-Inches(0.57), Inches(0.43), Inches(0.09), Inches(0.09))
    dot.fill.solid(); dot.fill.fore_color.rgb = RGBColor(0xE8,0x25,0x1A)
    dot.line.fill.background()

def _label(slide, text, accent):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(6), Inches(0.22))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = text.upper(); r.font.name = F_BODY; r.font.size = Pt(8)
    r.font.bold = True; r.font.color.rgb = accent

def _footer(slide):
    tb = slide.shapes.add_textbox(Inches(0.5), H-Inches(0.28), Inches(9), Inches(0.2))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = 'CONFIDENTIAL  \u00a9  20.20 Limited 2026'
    r.font.name = F_BODY; r.font.size = Pt(7.5); r.font.color.rgb = GREY

def _heading(slide, text, x=None, y=None, w=None, size=26, colour=None, caps=True):
    tb = slide.shapes.add_textbox(x or Inches(0.5), y or Inches(0.28),
                                   w or (W-Inches(1.2)), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text.upper() if caps else text
    r.font.name = F_HEAD; r.font.size = Pt(size)
    r.font.bold = True; r.font.color.rgb = colour or BLACK

def _textbox(slide, paras, x, y, w, h, size=12, colour=None):
    """paras: list of (type, text) — 'prose', 'bullet', 'bold'"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for typ, text in paras:
        if not text.strip(): continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if typ == 'bullet':
            from pptx.oxml.ns import qn
            from lxml import etree
            pPr = p._p.get_or_add_pPr()
            bc = etree.SubElement(pPr, qn('a:buChar'))
            bc.set('char', '\u2022')
            p.space_before = Pt(1)
        elif typ == 'prose':
            p.space_after = Pt(5)
        r = p.add_run(); r.text = text
        r.font.name = F_BODY; r.font.size = Pt(size)
        r.font.bold = (typ == 'bold')
        r.font.color.rgb = colour or MID

def _col_text(slide, items, x, y, w, h, size=11.5, heading=None, accent=None):
    """Column with optional coloured heading then bullets/prose."""
    cy = y
    if heading:
        tb = slide.shapes.add_textbox(x, cy, w, Inches(0.28))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = heading; r.font.name = F_HEAD; r.font.size = Pt(11)
        r.font.bold = True; r.font.color.rgb = accent or BLACK
        cy += Inches(0.3)
    _textbox(slide, items, x, cy, w, h - (Inches(0.3) if heading else 0), size)


# ── SLIDES ────────────────────────────────────────────────────────────────────

def slide_cover(prs, venue, client, contact, role, date_s, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, DARK)
    _box(slide, Inches(6.5), 0, Inches(6.83), H, RGBColor(0x1A,0x1E,0x24))
    tb = slide.shapes.add_textbox(Inches(7.0), Inches(0.2), Inches(5.8), H-Inches(0.4))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = '[IMAGE: Full-bleed stadium or hospitality photography — atmospheric, premium, on-brand]'
    r.font.name = F_BODY; r.font.size = Pt(9); r.font.color.rgb = RGBColor(0x44,0x44,0x44)

    # 20.20 mark
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(1.5), Inches(0.9))
    tf2 = tb2.text_frame
    for i, t in enumerate(['20','20']):
        p = tf2.paragraphs[0] if i==0 else tf2.add_paragraph()
        r2 = p.add_run(); r2.text = t
        r2.font.name = F_HEAD; r2.font.size = Pt(20)
        r2.font.bold = True; r2.font.color.rgb = WHITE

    _box(slide, Inches(0.5), Inches(1.45), Inches(0.75), Inches(0.05), accent)

    tb3 = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.7), Inches(1.4))
    tf3 = tb3.text_frame; tf3.word_wrap = True
    r3 = tf3.paragraphs[0].add_run()
    r3.text = venue; r3.font.name = F_HEAD; r3.font.size = Pt(38)
    r3.font.bold = True; r3.font.color.rgb = WHITE

    tb4 = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(5.7), Inches(0.4))
    r4 = tb4.text_frame.paragraphs[0].add_run()
    r4.text = 'Hospitality design proposal'
    r4.font.name = F_BODY; r4.font.size = Pt(14); r4.font.color.rgb = RGBColor(0xCC,0xCC,0xCC)

    tb5 = slide.shapes.add_textbox(Inches(0.5), H-Inches(1.3), Inches(5.7), Inches(0.85))
    tf5 = tb5.text_frame; tf5.word_wrap = True
    p5a = tf5.paragraphs[0]; r5a = p5a.add_run()
    r5a.text = client; r5a.font.name = F_HEAD; r5a.font.size = Pt(14)
    r5a.font.bold = True; r5a.font.color.rgb = WHITE
    p5b = tf5.add_paragraph(); r5b = p5b.add_run()
    r5b.text = f'Prepared for {contact}' + (f', {role}' if role else '') + (f'  |  {date_s}' if date_s else '')
    r5b.font.name = F_BODY; r5b.font.size = Pt(10); r5b.font.color.rgb = GREY


def slide_hello(prs, body, accent):
    """
    Cover letter slide — "Hello." headline, letter body, personal sign-off.
    Matches real 20.20 format: left-justified, personal, concise.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE); _logo(slide)

    # Project label top left (from brief metadata — added by caller)
    _label(slide, 'Introduction', accent)
    _rule(slide, Inches(0.5), Inches(0.88), W - Inches(1.0))

    # "Hello." headline — bold, left, with full stop
    # Extract "Hello [Name]" from first line of body if present
    lines_all = [l.strip() for l in body.splitlines() if l.strip()]
    hello_headline = 'Hello.'
    for line in lines_all[:2]:
        if line.lower().startswith('hello '):
            hello_headline = line.rstrip('.,')
            break

    tb_h = slide.shapes.add_textbox(Inches(0.5), Inches(0.94), Inches(9), Inches(0.62))
    p_hh = tb_h.text_frame.paragraphs[0]
    p_hh.alignment = PP_ALIGN.LEFT
    r_h = p_hh.add_run()
    r_h.text = hello_headline
    r_h.font.name = F_HEAD; r_h.font.size = Pt(30)
    r_h.font.bold = True; r_h.font.color.rgb = BLACK

    # Parse body — separate Dear/greeting, paragraphs, and sign-off
    clean_body = _clean(body)
    lines = [l.strip() for l in clean_body.splitlines() if l.strip()]

    # Find "Dear X" line
    greeting = ''
    sign_off_lines = []
    body_lines = []
    in_signoff = False

    for line in lines:
        if line.lower().startswith('dear '):
            greeting = line
        elif line.lower().startswith('kind regards') or line.lower().startswith('the 20.20'):
            in_signoff = True
            sign_off_lines.append(line)
        elif in_signoff:
            sign_off_lines.append(line)
        elif not greeting or body_lines or len(line) > 20:
            body_lines.append(line)

    # If no explicit sign-off found, add one
    if not sign_off_lines:
        sign_off_lines = ['Kind regards,', 'The 20.20 team']

    # Letter body — left third to two-thirds width (matching SAFC/CPFC style)
    letter_w = Inches(8.5)
    letter_x = Inches(0.5)
    y_start = Inches(1.65)

    # Greeting (Dear X,)
    if greeting:
        tb_g = slide.shapes.add_textbox(letter_x, y_start, letter_w, Inches(0.35))
        tf_g = tb_g.text_frame
        p_g = tf_g.paragraphs[0]
        p_g.alignment = PP_ALIGN.LEFT
        r_g = p_g.add_run()
        r_g.text = greeting
        r_g.font.name = F_BODY; r_g.font.size = Pt(12.5); r_g.font.color.rgb = MID
        y_start += Inches(0.48)

    # Body paragraphs
    paras_data = []
    for line in body_lines:
        if len(line) > 15:
            paras_data.append(('prose', line))

    if paras_data:
        _textbox(slide, paras_data, letter_x, y_start, letter_w,
                 H - y_start - Inches(1.2), size=12.5)

    # Sign-off block — bottom of letter
    sign_y = H - Inches(1.1)
    tb_s = slide.shapes.add_textbox(letter_x, sign_y, letter_w, Inches(0.8))
    tf_s = tb_s.text_frame
    first = True
    for sl in sign_off_lines:
        p = tf_s.paragraphs[0] if first else tf_s.add_paragraph()
        first = False
        r = p.add_run(); r.text = sl
        r.font.name = F_BODY; r.font.size = Pt(12)
        r.font.color.rgb = MID

    _footer(slide)


def slide_brief(prs, body, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE); _logo(slide); _label(slide, 'Your brief', accent)
    _rule(slide, Inches(0.5), Inches(0.92), W-Inches(1.0))
    _heading(slide, 'Our understanding', y=Inches(0.96), size=26)
    _rule(slide, Inches(0.5), Inches(1.68), W-Inches(1.0))

    all_items = [(t,c) for t,c in _parse(body)
                 if not _is_heading_line(c) and len(c) > 10]
    prose_items  = [(t,c) for t,c in all_items if t == 'prose']
    bullet_items = [(t,c) for t,c in all_items if t == 'bullet']

    # Always use two columns when we have enough content
    # Left: first 3 prose paragraphs
    # Right: remaining prose as bullets + any explicit bullets
    left_prose = prose_items[:3]
    right_prose = [('bullet', c) for t,c in prose_items[3:]]
    right_bullets = bullet_items  # use all bullets
    right_items = right_prose + right_bullets

    if right_items:
        _textbox(slide, left_prose, Inches(0.5), Inches(1.76),
                 Inches(5.9), Inches(5.3), size=12)
        _box(slide, Inches(6.55), Inches(1.76), Inches(0.01), Inches(5.0), RULE)
        _textbox(slide, right_items, Inches(6.7), Inches(1.76),
                 Inches(6.0), Inches(5.3), size=11.5)
    else:
        # Short brief — single column, max 4 paragraphs
        _textbox(slide, left_prose[:4], Inches(0.5), Inches(1.76),
                 W-Inches(1.0), Inches(5.3), size=12.5)
    _footer(slide)


def slide_process_summary(prs, stages, accent):
    """
    Overview slide showing all stages as numbered columns.
    stages: list of {'number', 'title', 'subtitle', 'riba', 'duration'}
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE); _logo(slide); _label(slide, 'Our methodology', accent)
    _rule(slide, Inches(0.5), Inches(0.92), W-Inches(1.0))
    _heading(slide, 'Our approach', y=Inches(0.96), size=26)

    # One-line intro
    tb_i = slide.shapes.add_textbox(Inches(0.5), Inches(1.68), W-Inches(1.0), Inches(0.38))
    r_i = tb_i.text_frame.paragraphs[0].add_run()
    r_i.text = 'Our process is structured across ' + str(len(stages)) + ' stages, each with clear deliverables and agreed review points.'
    r_i.font.name = F_BODY; r_i.font.size = Pt(12); r_i.font.color.rgb = MID

    # Dark RIBA bar
    _box(slide, Inches(0.5), Inches(2.12), W-Inches(1.0), Inches(0.32), BLACK)
    tb_bar = slide.shapes.add_textbox(Inches(0.6), Inches(2.12), Inches(6), Inches(0.32))
    r_bar = tb_bar.text_frame.paragraphs[0].add_run()
    stage_labels = list(dict.fromkeys(s.get('stage_label','') for s in stages if s.get('stage_label')))
    r_bar.text = '  |  '.join(stage_labels) if stage_labels else 'OUR PROCESS'
    r_bar.font.name = F_BODY; r_bar.font.size = Pt(9.5)
    r_bar.font.bold = True; r_bar.font.color.rgb = WHITE

    n = len(stages)
    col_w = (W - Inches(1.0)) / n
    top_y = Inches(2.55)

    for i, stage in enumerate(stages):
        x = Inches(0.5) + i * col_w
        # Large number
        tb_n = slide.shapes.add_textbox(x, top_y, Inches(0.65), Inches(0.8))
        r_n = tb_n.text_frame.paragraphs[0].add_run()
        r_n.text = str(stage.get('number', i+1))
        r_n.font.name = F_HEAD; r_n.font.size = Pt(40)
        r_n.font.bold = True; r_n.font.color.rgb = BLACK

        # Title
        tb_t = slide.shapes.add_textbox(x + Inches(0.65), top_y + Inches(0.1),
                                         col_w - Inches(0.75), Inches(0.65))
        tf_t = tb_t.text_frame; tf_t.word_wrap = True
        r_t = tf_t.paragraphs[0].add_run()
        r_t.text = stage.get('title','')
        r_t.font.name = F_HEAD; r_t.font.size = Pt(12)
        r_t.font.bold = True; r_t.font.color.rgb = BLACK

        # Subtitle / duration
        if stage.get('subtitle') or stage.get('duration'):
            tb_s = slide.shapes.add_textbox(x, top_y + Inches(0.82), col_w - Inches(0.1), Inches(0.28))
            r_s = tb_s.text_frame.paragraphs[0].add_run()
            r_s.text = stage.get('subtitle','') + (f"  |  {stage['duration']}" if stage.get('duration') else '')
            r_s.font.name = F_BODY; r_s.font.size = Pt(9.5)
            r_s.font.italic = True; r_s.font.color.rgb = GREY

        # Accent rule under number
        _box(slide, x, top_y + Inches(1.15), col_w - Inches(0.15), Inches(0.025), accent)

        # Key deliverables
        delivs = stage.get('deliverables', [])
        if delivs:
            _textbox(slide, [('bullet', d) for d in delivs[:5]],
                     x, top_y + Inches(1.22), col_w - Inches(0.15), Inches(3.8), size=10)

        # Column divider
        if i < n-1:
            _box(slide, x + col_w - Inches(0.06), top_y, Inches(0.01), Inches(4.5), RULE)

    _footer(slide)


def slide_stage_detail(prs, section_label, stage_title, body, accent):
    """
    Detailed stage slide: Objective | Process | Deliverables | Meetings
    Four columns matching real 20.20 format.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE); _logo(slide); _label(slide, section_label, accent)
    _rule(slide, Inches(0.5), Inches(0.92), W-Inches(1.0))
    _heading(slide, stage_title, y=Inches(0.96), size=24)

    # Intro line
    intro = _prose(body, 2)
    if intro:
        tb_i = slide.shapes.add_textbox(Inches(0.5), Inches(1.68), W-Inches(1.0), Inches(0.42))
        tf_i = tb_i.text_frame; tf_i.word_wrap = True
        r_i = tf_i.paragraphs[0].add_run()
        r_i.text = intro
        r_i.font.name = F_BODY; r_i.font.size = Pt(12); r_i.font.color.rgb = MID

    _rule(slide, Inches(0.5), Inches(2.16), W-Inches(1.0))

    # Extract structured sections — use ALL content, no caps
    obj_txt  = _section_of(body, 'Objective')  or _prose(body, 2)
    proc_txt = _section_of(body, 'Process')
    delv_txt = _section_of(body, 'Deliverables') or body
    meet_txt = _section_of(body, 'Meetings') or _section_of(body, 'Presentations')

    # Fall back: if no structured sections, split bullets across columns
    all_bullets = _all_bullets(body)
    if not proc_txt and not meet_txt:
        third = max(1, len(all_bullets) // 3)
        obj_bullets  = [obj_txt] if obj_txt and obj_txt != _prose(body,2) else all_bullets[:third]
        proc_bullets = all_bullets[third:third*2]
        delv_bullets = all_bullets[third*2:]
        meet_bullets = []
    else:
        obj_bullets  = [obj_txt] if obj_txt else []
        proc_bullets = _all_bullets(proc_txt) if proc_txt else []
        delv_bullets = _all_bullets(delv_txt) if delv_txt else all_bullets
        meet_bullets = _all_bullets(meet_txt) if meet_txt else []

    col_start = Inches(2.24)
    col_h     = H - col_start - Inches(0.35)
    col_w     = (W - Inches(1.0)) / 4
    # Default meetings cadence if not specified in content
    default_meetings = [
        'Kick-off meeting and site visit at start of stage',
        'Mid-stage review — design team and client',
        'End-of-stage presentation and sign-off',
        'All documents issued via PDF; meetings in person or on Teams',
    ]
    cols = [
        ('Objective',    [(t if t!='prose' else 'prose', c) for t,c in
                          [('prose',b) if isinstance(b,str) else b for b in obj_bullets]]),
        ('Process',      [('bullet',b) for b in proc_bullets]),
        ('Deliverables', [('bullet',b) for b in delv_bullets]),
        ('Meetings &\nPresentations', [('bullet',b) for b in meet_bullets] if meet_bullets
                          else [('bullet',b) for b in default_meetings]),
    ]

    for i, (col_label, items) in enumerate(cols):
        x = Inches(0.5) + i * col_w
        # Column heading
        tb_h = slide.shapes.add_textbox(x, col_start - Inches(0.06), col_w - Inches(0.1), Inches(0.3))
        r_h = tb_h.text_frame.paragraphs[0].add_run()
        r_h.text = col_label; r_h.font.name = F_HEAD; r_h.font.size = Pt(11)
        r_h.font.bold = True; r_h.font.color.rgb = accent
        # Divider
        if i < 3:
            _box(slide, x + col_w - Inches(0.06), col_start, Inches(0.01), col_h, RULE)
        _textbox(slide, items, x, col_start + Inches(0.28), col_w - Inches(0.1), col_h - Inches(0.28), size=10.5)

    _footer(slide)


def slide_fees(prs, stages_data, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE); _logo(slide); _label(slide, 'Our methodology', accent)
    _rule(slide, Inches(0.5), Inches(0.92), W-Inches(1.0))
    _heading(slide, 'Summary fees and timings', y=Inches(0.96), size=24)
    _rule(slide, Inches(0.5), Inches(1.68), W-Inches(1.0))

    # Column headers
    for txt, x in [('Fees',Inches(4.8)),('Timings',Inches(6.2)),('Invoicing schedule',Inches(7.5))]:
        tb = slide.shapes.add_textbox(x, Inches(1.76), Inches(3), Inches(0.25))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = txt; r.font.name = F_BODY; r.font.size = Pt(9); r.font.color.rgb = GREY

    row_h = Inches(0.88)
    y0 = Inches(2.06)

    for i, stage in enumerate(stages_data[:4]):
        y = y0 + i * row_h
        _box(slide, Inches(0.5), y, W-Inches(1.0), row_h-Inches(0.04), LGREY if i%2==0 else WHITE)
        # Number
        tb_n = slide.shapes.add_textbox(Inches(0.55), y+Inches(0.1), Inches(0.5), Inches(0.68))
        r_n = tb_n.text_frame.paragraphs[0].add_run()
        r_n.text = str(i+1); r_n.font.name = F_HEAD; r_n.font.size = Pt(28)
        r_n.font.bold = True; r_n.font.color.rgb = BLACK
        # Stage name
        tb_t = slide.shapes.add_textbox(Inches(1.1), y+Inches(0.1), Inches(3.55), Inches(0.68))
        tf_t = tb_t.text_frame; tf_t.word_wrap = True
        r_t = tf_t.paragraphs[0].add_run(); r_t.text = stage.get('title','').upper()
        r_t.font.name = F_BODY; r_t.font.size = Pt(10); r_t.font.bold = True; r_t.font.color.rgb = BLACK
        if stage.get('sub'):
            r_s = tf_t.add_paragraph().add_run(); r_s.text = stage['sub']
            r_s.font.name = F_BODY; r_s.font.size = Pt(9); r_s.font.color.rgb = MID
        # Fee
        tb_f = slide.shapes.add_textbox(Inches(4.7), y+Inches(0.2), Inches(1.4), Inches(0.45))
        r_f = tb_f.text_frame.paragraphs[0].add_run()
        r_f.text = stage.get('fee','[FEE: TBC]'); r_f.font.name = F_HEAD
        r_f.font.size = Pt(11); r_f.font.bold = True; r_f.font.color.rgb = BLACK
        # Timing
        tb_ti = slide.shapes.add_textbox(Inches(6.1), y+Inches(0.2), Inches(1.2), Inches(0.45))
        r_ti = tb_ti.text_frame.paragraphs[0].add_run()
        r_ti.text = stage.get('timing',''); r_ti.font.name = F_BODY
        r_ti.font.size = Pt(10); r_ti.font.color.rgb = MID
        # Invoicing
        tb_iv = slide.shapes.add_textbox(Inches(7.4), y+Inches(0.1), Inches(5.3), Inches(0.68))
        tb_iv.text_frame.word_wrap = True
        r_iv = tb_iv.text_frame.paragraphs[0].add_run()
        r_iv.text = stage.get('invoicing','50% at start, 50% at completion')
        r_iv.font.name = F_BODY; r_iv.font.size = Pt(9.5); r_iv.font.color.rgb = MID

    # Total
    tot_y = y0 + len(stages_data[:4]) * row_h + Inches(0.1)
    _rule(slide, Inches(0.5), tot_y, W-Inches(1.0), BLACK)
    tb_tt = slide.shapes.add_textbox(Inches(1.1), tot_y+Inches(0.1), Inches(3.5), Inches(0.4))
    r_tt = tb_tt.text_frame.paragraphs[0].add_run(); r_tt.text = 'TOTAL 20.20 FEES'
    r_tt.font.name = F_HEAD; r_tt.font.size = Pt(10); r_tt.font.bold = True; r_tt.font.color.rgb = BLACK
    tb_tf = slide.shapes.add_textbox(Inches(4.7), tot_y+Inches(0.1), Inches(1.4), Inches(0.35))
    r_tf = tb_tf.text_frame.paragraphs[0].add_run(); r_tf.text = '[FEE: TBC]'
    r_tf.font.name = F_HEAD; r_tf.font.size = Pt(11); r_tf.font.bold = True; r_tf.font.color.rgb = BLACK

    tb_sm = slide.shapes.add_textbox(Inches(0.5), H-Inches(0.48), Inches(9), Inches(0.26))
    r_sm = tb_sm.text_frame.paragraphs[0].add_run()
    r_sm.text = 'Fees are exclusive of VAT, 3rd party costs, general expenses and travel. Subject to contract.'
    r_sm.font.name = F_BODY; r_sm.font.size = Pt(8); r_sm.font.italic = True; r_sm.font.color.rgb = GREY
    _footer(slide)


def slide_divider(prs, word, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, DARK); _logo(slide, dark=True)
    tb = slide.shapes.add_textbox(Inches(0.5), H-Inches(2.2), W-Inches(1.0), Inches(1.8))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = word
    r.font.name = F_HEAD; r.font.size = Pt(72); r.font.bold = True; r.font.color.rgb = WHITE


def slide_next_steps(prs, body, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE); _logo(slide); _label(slide, 'Next steps', accent)
    _rule(slide, Inches(0.5), Inches(0.92), W-Inches(1.0))
    _heading(slide, 'Next steps', y=Inches(0.96), size=26)
    _rule(slide, Inches(0.5), Inches(1.68), W-Inches(1.0))

    bullets = _bullets(body, 4)
    defaults = [
        ('Review this proposal', 'Share with your team. Note any questions, changes or additions.'),
        ('Return your feedback', 'Send us your comments and we will revise and reissue.'),
        ('Site visit and kick-off', 'If you are minded to appoint us, we propose visiting the venue.'),
        ('Appointment', 'We can mobilise quickly from instruction.'),
    ]
    cw = Inches(2.9); gap = Inches(0.22)
    for i in range(4):
        x = Inches(0.5) + i*(cw+gap)
        _box(slide, x, Inches(1.85), cw, Inches(4.2), LGREY)
        tb_n = slide.shapes.add_textbox(x+Inches(0.18), Inches(2.0), Inches(0.6), Inches(0.52))
        r_n = tb_n.text_frame.paragraphs[0].add_run(); r_n.text = f'0{i+1}'
        r_n.font.name = F_HEAD; r_n.font.size = Pt(22); r_n.font.bold = True; r_n.font.color.rgb = accent
        _box(slide, x+Inches(0.18), Inches(2.56), Inches(1.5), Inches(0.025), accent)
        title = bullets[i] if i < len(bullets) else defaults[i][0]
        desc  = '' if i < len(bullets) else defaults[i][1]
        tb_t = slide.shapes.add_textbox(x+Inches(0.18), Inches(2.62), cw-Inches(0.3), Inches(0.55))
        tb_t.text_frame.word_wrap = True
        r_t = tb_t.text_frame.paragraphs[0].add_run(); r_t.text = title
        r_t.font.name = F_HEAD; r_t.font.size = Pt(11); r_t.font.bold = True; r_t.font.color.rgb = BLACK
        if desc:
            tb_d = slide.shapes.add_textbox(x+Inches(0.18), Inches(3.25), cw-Inches(0.3), Inches(2.3))
            tb_d.text_frame.word_wrap = True
            r_d = tb_d.text_frame.paragraphs[0].add_run(); r_d.text = desc
            r_d.font.name = F_BODY; r_d.font.size = Pt(10); r_d.font.color.rgb = MID
    _footer(slide)


# ── MAIN ──────────────────────────────────────────────────────────────────────

def find_section(sections, *keys):
    for key in keys:
        kl = key.lower()
        for sec in sections:
            h = sec.get('heading','').lower()
            if kl in h or h in kl:
                return _clean(sec.get('body',''))
    return ''

# Stage definitions by brief type
# RIBA-based proposals (newbuild, refurb, multi-stage)
STAGES_RIBA = [
    {'number':1,'title':'Strategic framework','subtitle':'Understand and define',
     'stage_label':'RIBA Stage 2','duration':'2-3 weeks',
     'deliverables':['Design principles per tier','Experience propositions','Naming and narratives','Design direction mood boards','Strategic report']},
    {'number':2,'title':'Concept design','subtitle':'Overarching look and feel',
     'stage_label':'RIBA Stage 2','duration':'4-6 weeks',
     'deliverables':['GA plans and zoning','Materials and mood boards','Brand identities','CGI visuals (min 2 per space)','Concept report']},
    {'number':3,'title':'Design development','subtitle':'Refine and finalise',
     'stage_label':'RIBA Stage 2','duration':'6-8 weeks',
     'deliverables':['Finalised GA plan and RCP','Sample boards','Furniture selection','Concept freeze','Design specification']},
    {'number':4,'title':'Design intent and artwork','subtitle':'Technical production',
     'stage_label':'RIBA Stage 3','duration':'8-12 weeks',
     'deliverables':['Full drawing pack','Elevations and sections','FFE schedules','Graphic artwork files','Specification document']},
    {'number':5,'title':'Coordination','subtitle':'On-site guardianship',
     'stage_label':'RIBA Stages 4-5','duration':'Programme dependent',
     'deliverables':['Contractor drawing review','Sample approvals','Site meetings','Value engineering support']},
    {'number':6,'title':'Handover','subtitle':'Completion and sign-off',
     'stage_label':'RIBA Stage 6','duration':'2 weeks',
     'deliverables':['Site inspection','Snagging report','Practical completion sign-off']},
]

# Phase-based proposals (arena, single space, sponsor, brand)
STAGES_PHASE = [
    {'number':1,'title':'Discovery and strategy','subtitle':'Understand and define',
     'stage_label':'Phase 1','duration':'2-3 weeks',
     'deliverables':['Project brief audit','Experience proposition','Design principles','Naming and narrative direction','Strategy document']},
    {'number':2,'title':'Concept design','subtitle':'Creative direction and look and feel',
     'stage_label':'Phase 2','duration':'4-6 weeks',
     'deliverables':['Space planning and layout','Materials and mood boards','Brand identity direction','CGI visuals','Concept report']},
    {'number':3,'title':'Design development','subtitle':'Refine, specify and freeze',
     'stage_label':'Phase 3','duration':'4-6 weeks',
     'deliverables':['Finalised design drawings','Sample boards','Furniture selection','Graphic artwork','Concept freeze document']},
    {'number':4,'title':'Production and delivery','subtitle':'Artwork and handover',
     'stage_label':'Phase 4','duration':'6-10 weeks',
     'deliverables':['Print-ready artwork files','Installation drawings','Supplier packs','Quality sign-off checklist']},
]

def get_stage_defs(brief_type, riba_stages):
    """Return appropriate stage definitions based on brief type."""
    bt = (brief_type or '').lower()
    riba = (riba_stages or '').upper()
    # Use phase-based for non-RIBA brief types
    if any(x in bt for x in ('arena','single','sponsor','brand','one space','lounge','suite')):
        return STAGES_PHASE, False
    # Use RIBA if explicitly mentioned
    if 'riba' in riba or any(x in bt for x in ('newbuild','refurb','multi','full')):
        return STAGES_RIBA, True
    # Default: RIBA for anything with multiple stages, phase for single
    return STAGES_RIBA, True

def build_pptx_clean(sections, meta, output_path):
    client  = meta.get('client','')
    venue   = meta.get('venue','Project')
    contact = meta.get('contact','')
    role    = meta.get('role','')
    date_s  = meta.get('date','')
    riba    = meta.get('riba_stages','')
    accent  = _accent(client)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # Work out which stages to include from brief type and riba_stages
    brief_type = meta.get('brief_type','')
    stage_defs, is_riba = get_stage_defs(brief_type, riba)

    stage_nums = []
    if riba:
        nums = re.findall(r'\d+', riba)
        if nums:
            lo, hi = int(nums[0]), int(nums[-1])
            stage_nums = list(range(lo, hi+1))
    if not stage_nums:
        stage_nums = [1,2,3,4]

    active_stages = [s for s in stage_defs if s['number'] in stage_nums]
    if not active_stages:
        active_stages = stage_defs[:4]

    # 1. Cover
    slide_cover(prs, venue, client, contact, role, date_s, accent)

    # 2. Hello / cover letter
    cover_txt = find_section(sections,'cover letter','cover','hello')
    # Prepend "Hello [Name]" if not already in text (fallback)
    if cover_txt and not cover_txt.lower().startswith('hello '):
        first_name = contact.split()[0] if contact else ''
        if first_name:
            cover_txt = f'Hello {first_name}\n\n' + cover_txt
    slide_hello(prs, cover_txt, accent)

    # 3. Brief
    brief_txt = find_section(sections,'your brief','brief reflection','understanding')
    slide_brief(prs, brief_txt, accent)

    # 4. Methodology divider
    slide_divider(prs, 'Our methodology', accent)

    # 5. Process summary
    slide_process_summary(prs, active_stages, accent)

    # 6-N. Stage detail slides — labels adapt to brief type and actual RIBA stage numbers
    continuation = meta.get('continuation','no').lower() == 'yes'
    prior_done = meta.get('prior_stages_completed','')

    if is_riba:
        # Build stage labels from actual stage_nums — respect continuation
        # e.g. if stage_nums = [3,4,5,6], label as Stage 3, Stage 4, etc.
        RIBA_LABELS = {
            1: ('stage 1','strategic framework', 'Stage 1 — Strategic framework'),
            2: ('stage 2','concept design',      'Stage 2 — Concept design'),
            3: ('stage 3','design development',  'Stage 3 — Design development'),
            4: ('stages 4','design intent',      'Stage 4 — Design intent and artwork'),
            5: ('stage 5','coordination',        'Stage 5 — Coordination and liaison'),
            6: ('stage 6','handover',            'Stage 6 — Handover'),
        }
        # Collapse 4+5+6 if all present
        working_nums = list(stage_nums)
        if 4 in working_nums and 5 in working_nums and 6 in working_nums:
            working_nums = [n for n in working_nums if n < 4] + [4]
            RIBA_LABELS[4] = ('stages 4','design intent', 'Stages 4, 5 and 6')

        stage_section_map = []
        for n in sorted(set(working_nums)):
            if n in RIBA_LABELS:
                k1, k2, label = RIBA_LABELS[n]
                stage_section_map.append((n, k1, k2, label))
    else:
        PHASE_LABELS = {
            1: ('stage 1','discovery',   'Phase 1 — Discovery and strategy'),
            2: ('stage 2','concept',     'Phase 2 — Concept design'),
            3: ('stage 3','development', 'Phase 3 — Design development'),
            4: ('stage 4','production',  'Phase 4 — Production and delivery'),
        }
        stage_section_map = [(n, *PHASE_LABELS[n]) for n in sorted(stage_nums) if n in PHASE_LABELS]

    for num, k1, k2, label in stage_section_map:
        txt = find_section(sections, k1, k2, 'stage '+str(num))
        if txt:
            slide_stage_detail(prs, 'Our methodology', label, txt, accent)

    # Fees
    fees_stages = [
        {'title':'Strategic framework','sub':'Workshop, site visit and proposition',
         'fee':'[FEE: TBC]','timing':'2-3 weeks','invoicing':'100% at start of stage'},
        {'title':'Concept design','sub':'Layouts, materials, CGI visuals',
         'fee':'[FEE: TBC]','timing':'4-6 weeks','invoicing':'50% at start, 50% at completion'},
        {'title':'Design development','sub':'Sample boards and concept freeze',
         'fee':'[FEE: TBC]','timing':'6-8 weeks','invoicing':'50% at start, 50% at completion'},
        {'title':'Design intent and artwork','sub':'Drawing pack and graphic artwork',
         'fee':'[FEE: TBC]','timing':'8-12 weeks','invoicing':'50% at start, 50% at completion'},
    ]
    slide_fees(prs, fees_stages, accent)

    # Next steps
    next_txt = find_section(sections,'next steps')
    slide_next_steps(prs, next_txt, accent)

    prs.save(output_path)
    return output_path
